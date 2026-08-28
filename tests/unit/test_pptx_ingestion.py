"""Lecture decks as knowledge files.

Every test builds a real .pptx with python-pptx and runs it through the real
parser. Nothing is stubbed: the things that broke on the first live deck — a
title ingested twice, bullets typed as prose, a footer ingested once per slide —
were all invisible to any test that asserted on a hand-built block list.

The shape being defended: a lecture slide is terse, and the sentences the
lecturer actually says live in the speaker notes. Ingesting only the slide keeps
the headings and throws away the teaching.
"""

import io

import pytest

from app.ai.ingestion.math_parser import (
    HEADING,
    LIST_ITEM,
    PARAGRAPH,
    TABLE,
    parse_blocks,
    parse_pptx_blocks,
)

pptx = pytest.importorskip("pptx", reason="python-pptx not installed")


def build(slides, *, footer=None):
    """Build a deck. `slides` is a list of (title, bullets, notes)."""
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    layout = presentation.slide_layouts[1]

    for title, bullets, notes in slides:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
        frame = slide.placeholders[1].text_frame
        for position, line in enumerate(bullets):
            paragraph = frame.paragraphs[0] if position == 0 else frame.add_paragraph()
            paragraph.text = line
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        if footer:
            box = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(3), Inches(0.4))
            box.text_frame.paragraphs[0].text = footer

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer


def texts(blocks, **match):
    return [
        b["text"] for b in blocks
        if all(b.get(k) == v for k, v in match.items())
    ]


DECK = [
    ("Photosynthesis",
     ["Light-dependent reactions", "Light-independent reactions"],
     "Photosynthesis converts light energy into chemical energy stored in glucose. "
     "Oxygen is a by-product, not the goal."),
    ("Respiration",
     ["Glycolysis", "Krebs cycle"],
     "Respiration releases the energy that photosynthesis stored."),
]


class TestSpeakerNotesAreIngested:
    """The reason this feature exists."""

    def test_notes_text_becomes_a_block(self):
        blocks = parse_pptx_blocks(build(DECK))
        notes = texts(blocks, is_speaker_notes=True)
        assert len(notes) == 2
        assert "chemical energy stored in glucose" in notes[0]

    def test_notes_are_attributed_to_their_slide(self):
        """A question about photosynthesis must retrieve the explanation AND be
        able to cite the slide it belongs to."""
        blocks = parse_pptx_blocks(build(DECK))
        note = next(b for b in blocks if b.get("is_speaker_notes"))
        assert note["heading_path"] == ["Photosynthesis"]
        assert note["page"] == 1

    def test_slide_content_is_not_marked_as_notes(self):
        blocks = parse_pptx_blocks(build(DECK))
        bullet = next(b for b in blocks if b["text"] == "Glycolysis")
        assert bullet["is_speaker_notes"] is False

    def test_a_deck_with_no_notes_still_parses(self):
        blocks = parse_pptx_blocks(build([("Title", ["A bullet"], None)]))
        assert texts(blocks, is_speaker_notes=True) == []
        assert "A bullet" in texts(blocks)


class TestSlideStructure:
    def test_the_title_becomes_a_heading_exactly_once(self):
        """`slide.shapes.title` returns a fresh proxy on every access, so an
        `is` check against the iterated shape never matches and the title is
        ingested as a heading AND again as body text."""
        blocks = parse_pptx_blocks(build(DECK))
        assert texts(blocks).count("Photosynthesis") == 1
        assert texts(blocks, type=HEADING) == ["Photosynthesis", "Respiration"]

    def test_bullets_are_list_items_not_paragraphs(self):
        """PowerPoint inherits the bullet glyph from the layout, so a real
        deck's bullets carry no marker of their own."""
        blocks = parse_pptx_blocks(build(DECK))
        assert "Glycolysis" in texts(blocks, type=LIST_ITEM)

    def test_the_slide_number_lands_in_page(self):
        """So a citation reads "p. 2 — Respiration", exactly as a PDF's would,
        and every downstream consumer of provenance keeps working."""
        blocks = parse_pptx_blocks(build(DECK))
        assert next(b for b in blocks if b["text"] == "Glycolysis")["page"] == 2

    def test_blocks_are_indexed_in_document_order(self):
        blocks = parse_pptx_blocks(build(DECK))
        assert [b["index"] for b in blocks] == list(range(len(blocks)))

    def test_a_table_survives_as_markdown(self):
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Limiting Factors"
        table = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(6), Inches(1)).table
        table.cell(0, 0).text = "Factor"
        table.cell(0, 1).text = "Effect"
        table.cell(1, 0).text = "Light"
        table.cell(1, 1).text = "Increases rate"
        buffer = io.BytesIO()
        presentation.save(buffer)
        buffer.seek(0)

        rendered = texts(parse_pptx_blocks(buffer), type=TABLE)
        assert len(rendered) == 1
        assert "| Factor | Effect |" in rendered[0]
        assert "| Light | Increases rate |" in rendered[0]


class TestFurnitureIsNotIngested:
    """A strap line repeated on every slide is not teaching material."""

    def test_a_repeated_footer_is_dropped(self):
        deck = [(f"Slide {n}", [f"Point {n}"], None) for n in range(1, 5)]
        blocks = parse_pptx_blocks(build(deck, footer="Dept of Biological Sciences"))
        assert not any("Biological Sciences" in t for t in texts(blocks))

    def test_content_repeated_only_twice_is_kept(self):
        """The threshold protects a phrase a lecturer genuinely repeats."""
        deck = [(f"Slide {n}", ["Remember the chain rule"], None) for n in range(1, 3)]
        blocks = parse_pptx_blocks(build(deck))
        assert texts(blocks).count("Remember the chain rule") == 2

    def test_a_repeated_slide_title_is_never_treated_as_furniture(self):
        """Decks legitimately reuse a section title across several slides."""
        deck = [("Worked Example", [f"Step {n}"], None) for n in range(1, 5)]
        blocks = parse_pptx_blocks(build(deck))
        assert texts(blocks, type=HEADING).count("Worked Example") == 4

    def test_repeated_notes_are_never_treated_as_furniture(self):
        deck = [(f"Slide {n}", ["x"], "Ask the class to try this one.") for n in range(1, 5)]
        blocks = parse_pptx_blocks(build(deck))
        assert len(texts(blocks, is_speaker_notes=True)) == 4


class TestTheDispatcherAcceptsDecks:
    def test_parse_blocks_routes_pptx(self):
        data = build(DECK).read()
        blocks = parse_blocks(data, "lecture-3.pptx")
        assert texts(blocks, type=HEADING) == ["Photosynthesis", "Respiration"]

    def test_extension_matching_is_case_insensitive(self):
        data = build(DECK).read()
        assert parse_blocks(data, "LECTURE-3.PPTX")

    def test_ppt_is_still_rejected_with_a_useful_message(self):
        """The old binary format is a different container; python-pptx cannot
        read it. Failing loudly beats a document stuck at 0 chunks."""
        with pytest.raises(ValueError) as excinfo:
            parse_blocks(b"anything", "legacy-deck.ppt")
        assert "pptx" in str(excinfo.value)


class TestTheUploadGateAgrees:
    """The parser supporting a format is not the same as the API accepting it —
    the allowlist is checked before ingestion ever runs."""

    def test_both_upload_paths_allow_pptx(self):
        from app.ai.ingestion.router import _ALLOWED_EXTENSIONS as ingestion_allowed
        from app.api.v1.lecturer.knowledge import _ALLOWED_EXTENSIONS as lecturer_allowed

        assert "pptx" in ingestion_allowed
        assert "pptx" in lecturer_allowed

    def test_every_allowed_extension_has_a_parser(self):
        """A file the API accepts but cannot parse is stored, queued, and then
        fails — the worst of both."""
        from app.ai.ingestion.math_parser import _PARSERS
        from app.api.v1.lecturer.knowledge import _ALLOWED_EXTENSIONS

        # .doc is the legacy binary Word format, handled upstream by conversion
        # rather than by a parser entry.
        unparseable = {e for e in _ALLOWED_EXTENSIONS if e not in _PARSERS} - {"doc"}
        assert not unparseable, f"accepted but unparseable: {sorted(unparseable)}"

    def test_pptx_uploads_get_the_right_content_type(self):
        """Storage serves the file back with whatever we stored it as."""
        from app.ai.ingestion.service import _content_type_for

        assert _content_type_for("lecture.pptx") == (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )


# ── Shapes PowerPoint hides from python-pptx ─────────────────────────────────
#
# These fixtures are built by hand rather than by python-pptx, because
# python-pptx never *writes* mc:AlternateContent — only PowerPoint does. That
# asymmetry is exactly why the bug below survived a green suite: a synthetic
# deck round-trips perfectly while a real one silently loses slides.

_MC = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_M = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_A14 = "http://schemas.microsoft.com/office/drawing/2010/main"

# One <m:oMath> holding x^2, wrapped the way PowerPoint wraps native equations:
# mc:Choice carries the real maths, mc:Fallback a flattened picture of it.
_ALTERNATE_CONTENT = f"""
<mc:AlternateContent xmlns:mc="{_MC}" xmlns:p="{_P}" xmlns:a="{_A}" xmlns:a14="{_A14}">
  <mc:Choice xmlns:m="{_M}" Requires="a14">
    <p:sp>
      <p:nvSpPr>
        <p:cNvPr id="99" name="Equation Box"/>
        <p:cNvSpPr txBox="1"/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="100" y="100"/><a:ext cx="500" cy="500"/></a:xfrm>
      </p:spPr>
      <p:txBody>
        <a:bodyPr/><a:lstStyle/>
        <a:p><a:r><a:t>Standard Limits</a:t></a:r></a:p>
        <a:p>
          <a14:m>
            <m:oMath>
              <m:sSup>
                <m:e><m:r><m:t>x</m:t></m:r></m:e>
                <m:sup><m:r><m:t>2</m:t></m:r></m:sup>
              </m:sSup>
            </m:oMath>
          </a14:m>
        </a:p>
      </p:txBody>
    </p:sp>
  </mc:Choice>
  <mc:Fallback>
    <p:sp>
      <p:nvSpPr>
        <p:cNvPr id="98" name="Equation Picture"/>
        <p:cNvSpPr txBox="1"/>
        <p:nvPr/>
      </p:nvSpPr>
      <p:spPr>
        <a:xfrm><a:off x="100" y="100"/><a:ext cx="500" cy="500"/></a:xfrm>
      </p:spPr>
      <p:txBody>
        <a:bodyPr/><a:lstStyle/>
        <a:p><a:r><a:t>FALLBACK PICTURE</a:t></a:r></a:p>
      </p:txBody>
    </p:sp>
  </mc:Fallback>
</mc:AlternateContent>
"""

# PowerPoint ink. python-pptx has no proxy class for <p:contentPart>, so lxml
# hands back a plain element and the shape factory dies on it. Found in a real
# lecturer deck.
_ALTERNATE_CONTENT_INK = f"""
<mc:AlternateContent xmlns:mc="{_MC}" xmlns:p="{_P}"
                     xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"
                     xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main">
  <mc:Choice Requires="p14">
    <p:contentPart r:id="rId99"/>
  </mc:Choice>
  <mc:Fallback/>
</mc:AlternateContent>
"""


def build_with_alternate_content(fragment=_ALTERNATE_CONTENT):
    """A one-slide deck carrying a hand-written mc:AlternateContent shape."""
    from lxml import etree
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Limits"
    slide.shapes._element.append(etree.fromstring(fragment))

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer


class TestAlternateContentIsNotInvisible:
    """PowerPoint wraps any shape using a post-2007 feature in
    mc:AlternateContent. Native equations are written that way, so on a maths
    deck this is the common case.

    python-pptx builds its shape tree from p:sp / p:grpSp / p:pic /
    p:graphicFrame / p:cxnSp only. The wrapper is none of those, so before the
    parser looked inside it, a slide whose equation *was* the content produced
    no blocks at all — no text, no maths, and no needs_math_ocr flag either,
    because that flag is set on a block and there was no block left to set it
    on. The upload still reported success.
    """

    def test_the_equation_inside_the_wrapper_is_converted(self):
        blocks = parse_pptx_blocks(build_with_alternate_content())
        equations = [e for b in blocks for e in (b.get("equations") or [])]
        assert equations, "the wrapped equation never reached the parser"
        assert any("x" in e and "2" in e for e in equations)

    def test_the_text_sharing_that_shape_is_kept(self):
        """The prose is the bigger loss: it goes with the equation."""
        blocks = parse_pptx_blocks(build_with_alternate_content())
        assert any("Standard Limits" in b["text"] for b in blocks)

    def test_the_slide_is_not_silently_empty(self):
        blocks = parse_pptx_blocks(build_with_alternate_content())
        assert len([b for b in blocks if b["type"] != HEADING]) > 0

    def test_the_choice_branch_wins_over_the_fallback(self):
        """Choice holds the maths as OMML, which converts exactly. Fallback
        holds a picture that could only be recovered by OCR — and taking both
        would ingest the slide twice."""
        blocks = parse_pptx_blocks(build_with_alternate_content())
        assert not any("FALLBACK PICTURE" in b["text"] for b in blocks)

    def test_an_unproxyable_shape_does_not_fail_the_deck(self):
        """<p:contentPart> is ink: nothing to extract, and no reason to lose
        the other 50 slides over it."""
        blocks = parse_pptx_blocks(build_with_alternate_content(_ALTERNATE_CONTENT_INK))
        assert any(b["text"] == "Limits" for b in blocks)


class TestTableCellsReportTheirMaths:
    def test_a_formula_in_a_table_cell_is_listed_in_equations(self):
        """The LaTeX was always inlined into the cell text, so retrieval worked.
        But the block claimed has_math with an empty equations list, which is a
        contradiction for any consumer reading the list rather than the text."""
        blocks = parse_pptx_blocks(build_with_alternate_content())
        for block in blocks:
            if block.get("has_math"):
                assert block["equations"], (
                    f"{block['type']} block claims maths but lists none"
                )


class TestFurnitureRemovalNeverEmptiesASlide:
    """The boilerplate fingerprint normalises digits to '#', which is right for
    the PDF furniture it was written for ("1 | P a g e") and wrong for a
    lecturer who titles a sequence "(2 of 8)" … "(8 of 8)" in a plain textbox
    rather than a title placeholder. Those collapse to one fingerprint, clear
    the three-slide threshold, and every slide in the series lost its only
    text — 7 slides across two real decks."""

    SERIES = [
        (f"Equation of a Straight Line ({n} of 8)", n) for n in range(1, 6)
    ]

    def _deck(self):
        """Titles in plain textboxes, as the real decks have them — a title
        placeholder would be exempt from furniture removal already."""
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        blank = presentation.slide_layouts[6]
        for title, _ in self.SERIES:
            slide = presentation.slides.add_slide(blank)
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
            box.text_frame.paragraphs[0].text = title
            footer = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(4), Inches(0.4))
            footer.text_frame.paragraphs[0].text = "INTRODUCTION TO CALCULUS"

        buffer = io.BytesIO()
        presentation.save(buffer)
        buffer.seek(0)
        return buffer

    def test_a_numbered_title_series_survives(self):
        blocks = parse_pptx_blocks(self._deck())
        kept = " ".join(b["text"] for b in blocks)
        for title, _ in self.SERIES:
            assert title in kept, f"{title!r} was dropped as furniture"

    def test_every_slide_still_produces_something(self):
        blocks = parse_pptx_blocks(self._deck())
        assert len({b["page"] for b in blocks}) == len(self.SERIES)

    def test_a_genuine_repeated_footer_is_still_dropped(self):
        """The rescue must not resurrect real furniture on slides that have
        other content to keep."""
        blocks = parse_pptx_blocks(self._deck())
        assert "INTRODUCTION TO CALCULUS" not in " ".join(b["text"] for b in blocks)
