"""Structure- and maths-preserving document parser.

Replaces the ``unstructured``-based parser for ingestion. Three things the old
parser lost, all of which matter for a Calculus / Statistics course:

1. **Equations.** ``python-docx``'s ``paragraph.text`` skips ``m:oMath``, so
   every Word equation was dropped. Here they are converted to LaTeX by
   :mod:`app.ai.ingestion.omml` and kept inline with the prose that explains
   them.
2. **Structure.** The old parser flattened a document into ``{title, content}``
   sections joined by single newlines, which left the chunker with nothing to
   split on. Here each block keeps its type, heading path, and document order.
3. **Provenance.** Nothing recorded where a passage came from, so citations were
   impossible. Every block carries its page (PDF) and heading path (both), which
   the chunker promotes into vector metadata.

Output is a flat list of :class:`Block` dicts in document order.
"""

from __future__ import annotations

import io
import logging
import os
import re
from typing import Any, Dict, Iterable, List, Optional

from app.ai.ingestion.omml import omml_to_latex

logger = logging.getLogger(__name__)

W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
# DrawingML: PowerPoint's text lives in this namespace, not Word's.
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Block types
HEADING = "heading"
PARAGRAPH = "paragraph"
EQUATION = "equation"
LIST_ITEM = "list_item"
TABLE = "table"
FIGURE = "figure"

Block = Dict[str, Any]

# Paragraph-level XML we must not walk into when collecting text: these hold
# formatting, not content.
_SKIP_TAGS = {
    "pPr", "rPr", "sectPr", "numPr", "tblPr", "tcPr", "trPr",
    "proofErr", "bookmarkEnd", "lastRenderedPageBreak", "commentRangeStart",
    "commentRangeEnd", "commentReference", "footnotePr", "endnotePr",
}


def _local(tag) -> str:
    if not isinstance(tag, str):
        return ""
    return tag.rsplit("}", 1)[-1]


def _escape_markdown_math(text: str) -> str:
    """Escape ``$`` so literal currency is not read as a math delimiter."""
    return text.replace("$", r"\$")


# ── DOCX ─────────────────────────────────────────────────────────────────────


def _collect_paragraph_parts(element, parts: List[tuple]) -> None:
    """Walk a ``w:p`` subtree in document order, keeping text and maths interleaved.

    Emits ``("text", str)``, ``("inline", latex)`` and ``("display", latex)``
    tuples. Order matters: "the derivative of $x^2$ is $2x$" must not collapse
    into "the derivative of is" with the maths appended somewhere else.
    """
    for child in element:
        tag = _local(child.tag)
        if not tag or tag in _SKIP_TAGS:
            continue

        if tag == "oMathPara":
            for equation in child.findall(f"{{{M_NS}}}oMath"):
                latex = omml_to_latex(equation)
                if latex:
                    parts.append(("display", latex))
        elif tag == "oMath":
            latex = omml_to_latex(child)
            if latex:
                parts.append(("inline", latex))
        elif tag == "t":
            parts.append(("text", child.text or ""))
        elif tag == "tab":
            parts.append(("text", " "))
        elif tag in ("br", "cr"):
            parts.append(("text", "\n"))
        elif tag == "drawing":
            parts.append(("figure", _drawing_description(child)))
        else:
            _collect_paragraph_parts(child, parts)


def _drawing_description(drawing) -> str:
    """Pull alt text off an embedded image so its position is not lost."""
    for prop in drawing.iter():
        if _local(prop.tag) in ("docPr", "cNvPr"):
            descr = prop.get("descr") or prop.get("name") or ""
            if descr:
                return descr.strip()
    return "image"


def _render_parts(parts: Iterable[tuple]) -> tuple[str, List[str], bool]:
    """Turn collected parts into markdown text plus the LaTeX found in them."""
    out: List[str] = []
    equations: List[str] = []
    display_only = True

    for kind, value in parts:
        if kind == "text":
            if value.strip():
                display_only = False
            out.append(_escape_markdown_math(value))
        elif kind == "inline":
            equations.append(value)
            out.append(f"${value}$")
        elif kind == "display":
            equations.append(value)
            out.append(f"\n\n$$\n{value}\n$$\n\n")
        elif kind == "figure":
            display_only = False
            out.append(f"![{value}](figure)")

    text = "".join(out)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return text, equations, display_only and bool(equations)


_HEADING_LEVEL_RE = re.compile(r"(\d+)")


def _heading_level(paragraph) -> Optional[int]:
    """Return the outline level of a paragraph, or None if it is body text."""
    try:
        style = paragraph.style.name or ""
    except Exception:
        style = ""

    if style.startswith("Heading"):
        match = _HEADING_LEVEL_RE.search(style)
        return int(match.group(1)) if match else 1
    if style in ("Title", "Subtitle"):
        return 1 if style == "Title" else 2

    # Some documents set an outline level without using a Heading style.
    outline = paragraph._p.find(f"{{{W_NS}}}pPr/{{{W_NS}}}outlineLvl")
    if outline is not None:
        val = outline.get(f"{{{W_NS}}}val")
        if val is not None and val.isdigit() and int(val) < 9:
            return int(val) + 1
    return None


def _is_list_item(paragraph) -> bool:
    if paragraph._p.find(f"{{{W_NS}}}pPr/{{{W_NS}}}numPr") is not None:
        return True
    # "List Bullet" / "List Number" carry their numbering on the style rather
    # than on the paragraph, so numPr alone misses them.
    try:
        return (paragraph.style.name or "").startswith("List")
    except Exception:
        return False


def _table_to_markdown(table) -> tuple[str, bool]:
    """Render a table as markdown, converting any equations inside its cells.

    The previous pipeline turned tables into prose ("normalize_sections"), which
    destroyed the row/column relationships that make a distribution table or a
    formula sheet useful.
    """
    rows: List[List[str]] = []
    has_math = False

    for row in table.rows:
        cells: List[str] = []
        for cell in row.cells:
            cell_parts: List[tuple] = []
            for paragraph in cell.paragraphs:
                _collect_paragraph_parts(paragraph._p, cell_parts)
                cell_parts.append(("text", " "))
            text, equations, _ = _render_parts(cell_parts)
            if equations:
                has_math = True
            cells.append(text.replace("\n", " ").replace("|", r"\|").strip())
        if any(cells):
            rows.append(cells)

    if not rows:
        return "", False

    return _rows_to_markdown(rows), has_math


def _rows_to_markdown(rows: List[List[str]]) -> str:
    """Render already-extracted cells as a markdown table.

    Split out of _table_to_markdown so the PPTX parser can reuse it. A Word
    table and a PowerPoint table expose different objects (cell.paragraphs
    vs cell.text_frame.paragraphs) but yield the same rows, and a formula
    sheet should render identically whichever file it arrived in.
    """
    if not rows:
        return ""

    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]

    header, *body = rows
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * width) + " |",
    ]
    lines += ["| " + " | ".join(r) + " |" for r in body]
    return "\n".join(lines)


def parse_docx_blocks(source) -> List[Block]:
    """Parse a .docx into ordered blocks, preserving equations and structure."""
    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    document = Document(source)
    body = document.element.body

    blocks: List[Block] = []
    heading_path: List[str] = []
    index = 0

    for element in body:
        tag = _local(element.tag)

        if tag == "p":
            paragraph = Paragraph(element, document)
            parts: List[tuple] = []
            _collect_paragraph_parts(element, parts)
            text, equations, equation_only = _render_parts(parts)
            if not text:
                continue

            level = _heading_level(paragraph)
            if level is not None:
                heading_path = heading_path[: level - 1]
                heading_path.append(text)
                blocks.append({
                    "type": HEADING,
                    "text": text,
                    "level": level,
                    "heading_path": list(heading_path[:-1]),
                    "page": None,
                    "equations": equations,
                    "has_math": bool(equations),
                    "index": index,
                })
            else:
                if equation_only:
                    block_type = EQUATION
                    # Word marks a standalone equation with a bare m:oMath as
                    # often as with m:oMathPara. A paragraph that is nothing but
                    # maths is display maths either way, so promote it rather
                    # than leaving it inline mid-line.
                    text = "$$\n" + "\n".join(equations) + "\n$$"
                elif _is_list_item(paragraph):
                    block_type = LIST_ITEM
                else:
                    block_type = PARAGRAPH
                blocks.append({
                    "type": block_type,
                    "text": text,
                    "level": None,
                    "heading_path": list(heading_path),
                    "page": None,
                    "equations": equations,
                    "has_math": bool(equations),
                    "index": index,
                })
            index += 1

        elif tag == "tbl":
            markdown, has_math = _table_to_markdown(Table(element, document))
            if not markdown:
                continue
            blocks.append({
                "type": TABLE,
                "text": markdown,
                "level": None,
                "heading_path": list(heading_path),
                "page": None,
                "equations": [],
                "has_math": has_math,
                "index": index,
            })
            index += 1

    return blocks


# ── PDF ──────────────────────────────────────────────────────────────────────

# A PDF text layer stores glyphs, not semantics, so an equation typeset by LaTeX
# comes out as mangled runs like "f (x) = x2 1". We keep the text (it is still
# searchable prose) and flag the block so downstream code can warn, rather than
# pretending the maths survived. Set MATH_OCR_PROVIDER to route these through a
# real formula OCR service when one is configured.
_MATH_GLYPH_RE = re.compile(r"[∫∑∏√±≤≥≠≈∞∂∇αβγδθλμσπΣΔΩ]")
# Below this many characters a page has no usable text layer — treat it as a scan.
_MIN_PAGE_TEXT_CHARS = 40
_EQUATION_HINT_RE = re.compile(r"[a-zA-Z]\s*[=<>]\s*[-+]?[\w\\(]|\bd[xyt]\b|\^|_\{")


def _looks_mathematical(text: str) -> bool:
    return bool(_MATH_GLYPH_RE.search(text) or _EQUATION_HINT_RE.search(text))


_PDF_HEADING_RE = re.compile(
    r"^(?:\d+(?:\.\d+)*\.?\s+\S|CHAPTER\b|SECTION\b|Chapter\b|Section\b)"
)


def _looks_like_heading(line: str) -> bool:
    stripped = line.strip()
    if not stripped or len(stripped) > 90:
        return False
    if _PDF_HEADING_RE.match(stripped):
        return True
    # A short line in title case or caps with no terminal punctuation.
    if len(stripped.split()) <= 10 and not stripped.endswith((".", ",", ";", ":")):
        letters = [c for c in stripped if c.isalpha()]
        if letters and sum(c.isupper() for c in letters) / len(letters) > 0.6:
            return True
    return False


def _boilerplate_fingerprint(line: str) -> str:
    """Normalise a line so running headers/footers match across pages.

    Page furniture varies only by number and is often letter-spaced by the PDF
    text layer ("1 | P a g e"), so digits are collapsed to ``#`` and all
    whitespace is removed before comparing.
    """
    return re.sub(r"\s+", "", re.sub(r"\d+", "#", line)).lower()


def _verbatim_fingerprint(line: str) -> str:
    """As above but keeping digits, so "(2 of 8)" and "(3 of 8)" stay distinct.

    Used where collapsing the numbers would merge a numbered *title series*
    into one running footer — see :func:`_pptx_drop_repeated_furniture`.
    """
    return re.sub(r"\s+", "", line).lower()


def _find_boilerplate(page_texts: List[str]) -> set:
    """Identify running headers and footers so they do not become headings."""
    from collections import Counter

    if len(page_texts) < 2:
        return set()

    counts: Counter = Counter()
    for text in page_texts:
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        # On a sparse page every line is content — there is no furniture to
        # separate from the body, and guessing costs us real material.
        if len(lines) < 4:
            continue
        # Furniture sits on the outermost line at either edge, and is short.
        for line in (lines[0], lines[-1]):
            if len(line) <= 80:
                counts[_boilerplate_fingerprint(line)] += 1

    threshold = max(2, int(len(page_texts) * 0.3))
    return {fp for fp, n in counts.items() if n >= threshold and fp}


def parse_pdf_blocks(source) -> List[Block]:
    """Parse a PDF into ordered blocks, one group per paragraph, page numbers kept."""
    from pypdf import PdfReader

    reader = PdfReader(source)
    blocks: List[Block] = []
    heading_path: List[str] = []
    index = 0

    page_texts: List[str] = []
    for page_number, page in enumerate(reader.pages, start=1):
        try:
            page_texts.append(page.extract_text() or "")
        except Exception:
            logger.warning("Failed to extract text from page %s", page_number, exc_info=True)
            page_texts.append("")

    boilerplate = _find_boilerplate(page_texts)

    for page_number, raw in enumerate(page_texts, start=1):
        # A page with (almost) no text layer is a scan or an image-only export —
        # a lecturer's photographed or scanned handwritten notes look exactly
        # like this. Without a marker the page would vanish silently, so emit a
        # placeholder that routes it to OCR instead.
        if len(raw.strip()) < _MIN_PAGE_TEXT_CHARS:
            blocks.append({
                "type": PARAGRAPH,
                "text": "",
                "level": None,
                "heading_path": list(heading_path),
                "page": page_number,
                "equations": [],
                "has_math": False,
                "needs_math_ocr": True,
                "is_image_page": True,
                "index": index,
            })
            index += 1
            continue

        raw = "\n".join(
            ln for ln in raw.splitlines()
            if _boilerplate_fingerprint(ln) not in boilerplate
        )
        # Join hyphenated line breaks, then group into paragraphs on blank lines.
        raw = re.sub(r"-\n(?=[a-z])", "", raw)
        for chunk in re.split(r"\n\s*\n", raw):
            lines = [ln.rstrip() for ln in chunk.splitlines() if ln.strip()]
            if not lines:
                continue

            if len(lines) == 1 and _looks_like_heading(lines[0]):
                heading = lines[0].strip()
                heading_path = [heading]
                blocks.append({
                    "type": HEADING,
                    "text": heading,
                    "level": 1,
                    "heading_path": [],
                    "page": page_number,
                    "equations": [],
                    "has_math": False,
                    "index": index,
                })
                index += 1
                continue

            text = " ".join(lines).strip()
            if not text:
                continue
            blocks.append({
                "type": PARAGRAPH,
                "text": _escape_markdown_math(text),
                "level": None,
                "heading_path": list(heading_path),
                "page": page_number,
                "equations": [],
                "has_math": False,
                "needs_math_ocr": _looks_mathematical(text),
                "index": index,
            })
            index += 1

    return blocks



# ── PPTX ─────────────────────────────────────────────────────────────────────

# Placeholders that carry furniture rather than teaching. Without this the slide
# number, the date and the department strap line are ingested once per slide,
# and a 60-slide deck contributes 60 chunks that say "3".
_PPTX_FURNITURE = {"SLIDE_NUMBER", "DATE", "FOOTER"}

# Slides are terser than prose: a bullet is often three words. Below this a
# block is kept but marked, so the chunker can glue it to its neighbours rather
# than embedding "Therefore:" as a standalone passage.
_PPTX_MIN_STANDALONE_CHARS = 25


# Markup Compatibility namespace. Not in python-pptx's own prefix map, so `qn`
# raises KeyError on "mc:…" and the URI has to be written out.
_MC_NS = "{http://schemas.openxmlformats.org/markup-compatibility/2006}"
_P_NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

# The shape elements python-pptx can build a proxy for. Anything else inside an
# mc:Choice is skipped by name rather than by catching the failure: a real deck
# put a <p:contentPart> (PowerPoint ink) in one, lxml returned a plain _Element
# because python-pptx has no class for it, and the factory died reaching for an
# attribute only its own element classes have. Ink carries no text anyway, so
# there is nothing to recover — but the deck must not fail over it.
_PPTX_SHAPE_TAGS = frozenset(
    f"{_P_NS}{tag}" for tag in ("sp", "grpSp", "pic", "graphicFrame", "cxnSp")
)


def _pptx_recover_alternate_content(shapes) -> List[tuple]:
    """Recover the shapes PowerPoint hides inside ``<mc:AlternateContent>``.

    PowerPoint wraps any shape using a feature an older reader cannot render in
    ``mc:AlternateContent``: ``<mc:Choice Requires="a14">`` carries the real
    shape and ``<mc:Fallback>`` a flattened picture of it. **Native equations
    are written this way**, so on a mathematics deck this is the common case
    rather than an edge case.

    python-pptx builds its shape tree only from ``p:sp``, ``p:grpSp``, ``p:pic``,
    ``p:graphicFrame`` and ``p:cxnSp``. The wrapper is none of those, so it is
    not a shape to python-pptx and everything inside it is invisible. On a slide
    whose equation *is* the content that yields no blocks at all — no text, no
    maths, and no ``needs_math_ocr`` flag either, because that flag is set on a
    block and there is no block left to set it on. The slide's prose goes with
    it, silently, and the upload still reports success.

    ``Choice`` is preferred over ``Fallback`` deliberately: it holds the maths
    as OMML, which :func:`omml_to_latex` converts exactly, whereas ``Fallback``
    holds a picture that could only be recovered by OCR.

    Returns ``(document_position, shape)`` pairs so the caller can splice them
    back in where they belong rather than appending them after the slide.
    """
    parent = getattr(shapes, "_element", None)
    factory = getattr(shapes, "_shape_factory", None)
    if parent is None or factory is None:
        return []

    recovered: List[tuple] = []
    for position, child in enumerate(parent):
        if child.tag != f"{_MC_NS}AlternateContent":
            continue

        branch = child.find(f"{_MC_NS}Choice")
        if branch is None:
            branch = child.find(f"{_MC_NS}Fallback")
        if branch is None:
            continue

        for element in branch:
            if element.tag not in _PPTX_SHAPE_TAGS:
                logger.debug(
                    "Skipping <%s> inside mc:AlternateContent — not a shape element",
                    _local(element.tag),
                )
                continue
            try:
                recovered.append((position, factory(element)))
            except Exception as exc:  # noqa: BLE001
                # Backstop for a shape element this python-pptx cannot proxy.
                # Losing one shape is the right cost; raising would lose the
                # whole deck over a construct we do not need to understand.
                logger.warning(
                    "Skipped <%s> inside mc:AlternateContent: %s",
                    _local(element.tag), exc,
                )
    return recovered


def _pptx_shapes_in_document_order(shapes) -> List[Any]:
    """python-pptx's shapes plus the ones it cannot see, in document order."""
    parent = getattr(shapes, "_element", None)
    if parent is None:
        return list(shapes)

    order = {id(element): position for position, element in enumerate(parent)}
    end = len(order)

    numbered: List[tuple] = [
        (order.get(id(shape._element), end), shape) for shape in shapes
    ]
    numbered.extend(_pptx_recover_alternate_content(shapes))

    # Stable, so several shapes sharing one wrapper keep their branch order.
    numbered.sort(key=lambda item: item[0])
    return [shape for _, shape in numbered]


def _pptx_iter_shapes(shapes) -> Iterable[Any]:
    """Yield shapes in reading order, descending into groups.

    PowerPoint stores shapes in z-order, which is the order they were drawn in,
    not the order they are read in. A lecturer who adds a heading last would
    otherwise have it ingested last, after the bullets it introduces. Sorting
    top-to-bottom then left-to-right recovers the reading order; shapes with no
    geometry keep their document position.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    flat: List[Any] = []
    for shape in _pptx_shapes_in_document_order(shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            flat.extend(_pptx_iter_shapes(shape.shapes))
        else:
            flat.append(shape)

    def position(item) -> tuple:
        top = item[1].top
        left = item[1].left
        # None means "inherits from the layout" — there is no honest coordinate
        # to sort on, so fall back to document order for that shape alone.
        if top is None or left is None:
            return (1, 0, 0, item[0])
        return (0, top, left, item[0])

    return [shape for _, shape in sorted(enumerate(flat), key=position)]


def _pptx_is_furniture(shape) -> bool:
    """True for slide-number / date / footer placeholders."""
    try:
        if not shape.is_placeholder:
            return False
        return str(shape.placeholder_format.type).split()[0].upper() in _PPTX_FURNITURE
    except (AttributeError, ValueError):
        return False


def _pptx_is_body_placeholder(shape) -> bool:
    """Whether a shape is the slide's content area rather than a loose textbox."""
    try:
        if not shape.is_placeholder:
            return False
        kind = str(shape.placeholder_format.type).split()[0].upper()
    except (AttributeError, ValueError):
        return False
    return kind in {"BODY", "OBJECT", "SUBTITLE", "CONTENT"}


def _pptx_is_bullet(paragraph, *, in_body_placeholder: bool) -> bool:
    """Whether a slide paragraph renders with a bullet.

    An indented paragraph is always a bullet. Otherwise PowerPoint inherits the
    glyph from the layout, so a normal deck's bullets carry no marker of their
    own — testing only for an explicit ``a:buChar`` types every bullet on every
    slide as a plain paragraph. Inheritance applies to body placeholders; a
    free-floating textbox inherits nothing.
    """
    if paragraph.level:
        return True

    properties = paragraph._p.find(f"{{{_A_NS}}}pPr")
    if properties is not None:
        if properties.find(f"{{{_A_NS}}}buNone") is not None:
            return False
        if (properties.find(f"{{{_A_NS}}}buChar") is not None
                or properties.find(f"{{{_A_NS}}}buAutoNum") is not None):
            return True

    return in_body_placeholder


def _pptx_text_blocks(
    text_frame,
    *,
    slide_number: int,
    heading_path: List[str],
    start_index: int,
    is_notes: bool = False,
    in_body_placeholder: bool = False,
) -> List[Block]:
    """Turn one text frame into blocks, keeping any equations it contains."""
    blocks: List[Block] = []
    index = start_index

    for paragraph in text_frame.paragraphs:
        parts: List[tuple] = []
        # Reused from the DOCX path on purpose: it dispatches on local tag
        # names, and DrawingML runs are <a:t> — the same local name as Word's
        # <w:t>. Equations embedded in a slide (<a14:m><m:oMath>) are reached by
        # its recursion and converted to LaTeX, so maths on a slide survives
        # exactly as it does in a Word document.
        _collect_paragraph_parts(paragraph._p, parts)
        text, equations, equation_only = _render_parts(parts)
        if not text:
            continue

        if equation_only:
            block_type = EQUATION
            text = "$$\n" + "\n".join(equations) + "\n$$"
        elif _pptx_is_bullet(paragraph, in_body_placeholder=in_body_placeholder):
            block_type = LIST_ITEM
        else:
            block_type = PARAGRAPH

        blocks.append({
            "type": block_type,
            "text": text,
            "level": None,
            "heading_path": list(heading_path),
            "page": slide_number,
            "equations": equations,
            "has_math": bool(equations),
            # A slide can hold its maths as a pasted image, which python-pptx
            # cannot read. Flagging it routes the slide to the same OCR pass a
            # scanned PDF page gets.
            "needs_math_ocr": not equations and _looks_mathematical(text),
            "is_speaker_notes": is_notes,
            "is_terse": len(text) < _PPTX_MIN_STANDALONE_CHARS,
            "index": index,
        })
        index += 1

    return blocks


def _pptx_drop_repeated_furniture(blocks: List[Block]) -> List[Block]:
    """Remove the course code / department strap repeated on every slide.

    Lecturers put these in a plain textbox as often as in a footer placeholder,
    where nothing marks them as furniture. What gives them away is repetition:
    a short line appearing verbatim on several slides is a running footer, not
    teaching. Same fingerprint the PDF parser uses for running headers.

    Headings and speaker notes are exempt. A deck can legitimately reuse a
    section title, and notes are never furniture.

    Two strengths of match, because they carry different confidence:

    * **verbatim** — the same line, digits and all, on three or more slides.
      That is a running footer and is always dropped.
    * **digit-normalised only** — lines that match *after* numbers collapse to
      ``#``. That rule exists for the PDF furniture this fingerprint was
      written for ("1 | P a g e"), and on a deck it also catches a lecturer who
      titles a sequence "Equation of a Straight Line (2 of 8)" … "(8 of 8)" in
      a plain textbox rather than a title placeholder. Dropped too — but never
      as the last block on its slide, because a running footer is never the
      only thing on a slide, whereas one of those titles is.

    Without that second distinction the series above cleared the three-slide
    threshold and every slide in it lost its only text: 7 slides across two
    real decks, silently.
    """
    def considered(block) -> bool:
        return (
            block["type"] != HEADING
            and not block.get("is_speaker_notes")
            and len(block["text"]) <= 80
        )

    verbatim_seen: Dict[str, set] = {}
    collapsed_seen: Dict[str, set] = {}
    for block in blocks:
        if not considered(block):
            continue
        verbatim_seen.setdefault(
            _verbatim_fingerprint(block["text"]), set()
        ).add(block["page"])
        collapsed_seen.setdefault(
            _boilerplate_fingerprint(block["text"]), set()
        ).add(block["page"])

    verbatim = {fp for fp, pages in verbatim_seen.items() if len(pages) >= 3}
    collapsed = {fp for fp, pages in collapsed_seen.items() if len(pages) >= 3}
    if not verbatim and not collapsed:
        return blocks

    def is_furniture(block) -> bool:
        if not considered(block):
            return False
        return (
            _verbatim_fingerprint(block["text"]) in verbatim
            or _boilerplate_fingerprint(block["text"]) in collapsed
        )

    def is_verbatim_furniture(block) -> bool:
        return (
            considered(block)
            and _verbatim_fingerprint(block["text"]) in verbatim
        )

    kept = [b for b in blocks if not is_furniture(b)]

    # Restore a slide the pass emptied, but only from the weaker matches: a
    # verbatim running footer stays dropped even when it is all a slide has.
    surviving = {b["page"] for b in kept}
    if len(surviving) < len({b["page"] for b in blocks}):
        keep_ids = {id(b) for b in kept}
        kept = [
            b for b in blocks
            if id(b) in keep_ids
            or (b["page"] not in surviving and not is_verbatim_furniture(b))
        ]

    for position, block in enumerate(kept):
        block["index"] = position
    return kept


def parse_pptx_blocks(source) -> List[Block]:
    """Parse a .pptx into ordered blocks, one group per slide.

    The slide number goes in ``page``, so a citation reads "p. 4 — Photosynthesis"
    exactly as a PDF's would, and every downstream consumer of provenance keeps
    working without knowing decks exist.

    **Speaker notes are ingested.** On a real lecture deck the slide is five
    bullet fragments and the notes are the sentences the lecturer actually says;
    ingesting only the slide would keep the headings and throw away the
    teaching. Notes are attributed to their slide's heading, so a question about
    photosynthesis retrieves the explanation and cites the slide it belongs to.
    """
    from pptx import Presentation

    presentation = Presentation(source)
    blocks: List[Block] = []
    index = 0

    for slide_number, slide in enumerate(presentation.slides, start=1):
        heading_path: List[str] = []

        title_shape = None
        try:
            title_shape = slide.shapes.title
        except (AttributeError, ValueError):
            title_shape = None

        if title_shape is not None and title_shape.has_text_frame:
            title = " ".join(title_shape.text_frame.text.split()).strip()
            if title:
                heading_path = [title]
                blocks.append({
                    "type": HEADING,
                    "text": title,
                    "level": 1,
                    "heading_path": [],
                    "page": slide_number,
                    "equations": [],
                    "has_math": False,
                    "index": index,
                })
                index += 1

        for shape in _pptx_iter_shapes(slide.shapes):
            # Identity, not `is`: python-pptx returns a new proxy object on
            # every `shapes.title` access, so the obvious check silently
            # never matches and every title is ingested twice.
            same_as_title = (
                title_shape is not None and shape._element is title_shape._element
            )
            if same_as_title or _pptx_is_furniture(shape):
                continue

            if getattr(shape, "has_table", False):
                # Cells go through the same part collector as everything else,
                # so a formula sheet pasted into a slide table keeps its maths
                # instead of arriving as the empty string.
                rows: List[List[str]] = []
                # Kept, not just counted. The LaTeX is already inlined into the
                # cell text, so retrieval was never broken — but every other
                # block reports its maths here, and a block claiming has_math
                # with an empty equations list makes the two disagree for any
                # consumer that reads the list rather than the text.
                table_equations: List[str] = []
                for row in shape.table.rows:
                    cells: List[str] = []
                    for cell in row.cells:
                        cell_parts: List[tuple] = []
                        for paragraph in cell.text_frame.paragraphs:
                            _collect_paragraph_parts(paragraph._p, cell_parts)
                            cell_parts.append(("text", " "))
                        cell_text, cell_equations, _ = _render_parts(cell_parts)
                        table_equations.extend(cell_equations)
                        cells.append(
                            cell_text.replace("\n", " ").replace("|", r"\|").strip()
                        )
                    if any(cells):
                        rows.append(cells)

                text = _rows_to_markdown(rows)
                if text:
                    blocks.append({
                        "type": TABLE,
                        "text": text,
                        "level": None,
                        "heading_path": list(heading_path),
                        "page": slide_number,
                        "equations": table_equations,
                        "has_math": bool(table_equations),
                        "index": index,
                    })
                    index += 1
                continue

            if getattr(shape, "has_text_frame", False):
                produced = _pptx_text_blocks(
                    shape.text_frame,
                    slide_number=slide_number,
                    heading_path=heading_path,
                    start_index=index,
                    in_body_placeholder=_pptx_is_body_placeholder(shape),
                )
                blocks.extend(produced)
                index += len(produced)

        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame is not None and notes_frame.text.strip():
                produced = _pptx_text_blocks(
                    notes_frame,
                    slide_number=slide_number,
                    heading_path=heading_path,
                    start_index=index,
                    is_notes=True,
                )
                blocks.extend(produced)
                index += len(produced)

    return _pptx_drop_repeated_furniture(blocks)


# ── Plain text / markdown ────────────────────────────────────────────────────

_MD_HEADING_RE = re.compile(r"^(#{1,6})\s+(.*)$")
_MD_MATH_RE = re.compile(r"\$\$(.+?)\$\$|\$([^$\n]+)\$", re.DOTALL)


def parse_text_blocks(source) -> List[Block]:
    """Parse .md / .txt. LaTeX written as $…$ or $$…$$ is kept as-is."""
    data = source.read() if hasattr(source, "read") else source
    if isinstance(data, bytes):
        data = data.decode("utf-8", errors="replace")

    blocks: List[Block] = []
    heading_path: List[str] = []
    index = 0

    for chunk in re.split(r"\n\s*\n", data):
        text = chunk.strip()
        if not text:
            continue

        heading = _MD_HEADING_RE.match(text)
        if heading and "\n" not in text:
            level = len(heading.group(1))
            title = heading.group(2).strip()
            heading_path = heading_path[: level - 1]
            heading_path.append(title)
            blocks.append({
                "type": HEADING, "text": title, "level": level,
                "heading_path": list(heading_path[:-1]), "page": None,
                "equations": [], "has_math": False, "index": index,
            })
            index += 1
            continue

        equations = [m.group(1) or m.group(2) for m in _MD_MATH_RE.finditer(text)]
        blocks.append({
            "type": PARAGRAPH, "text": text, "level": None,
            "heading_path": list(heading_path), "page": None,
            "equations": [e.strip() for e in equations if e],
            "has_math": bool(equations), "index": index,
        })
        index += 1

    return blocks


# ── Entry point ──────────────────────────────────────────────────────────────

_PARSERS = {
    "docx": parse_docx_blocks,
    "docm": parse_docx_blocks,
    "pdf": parse_pdf_blocks,
    "pptx": parse_pptx_blocks,
    "md": parse_text_blocks,
    "markdown": parse_text_blocks,
    "txt": parse_text_blocks,
}

SUPPORTED_EXTENSIONS = tuple(sorted(_PARSERS))


def parse_blocks(content: bytes, filename: str) -> List[Block]:
    """Parse raw file bytes into ordered blocks.

    Raises ``ValueError`` for an unsupported extension, and for bytes that are
    not the format the extension claims, so ingestion can mark the document
    failed with a message the admin UI can show.
    """
    suffix = os.path.splitext(filename)[1].lstrip(".").lower()
    parser = _PARSERS.get(suffix)
    if parser is None:
        raise ValueError(
            f"Unsupported file type {suffix!r}. Supported: {', '.join(SUPPORTED_EXTENSIONS)}"
        )

    try:
        blocks = parser(io.BytesIO(content))
    except ValueError:
        raise
    except Exception as exc:
        # The parser choked on the container itself, which means the bytes are
        # not the file the extension claims: a truncated upload, an HTML error
        # page saved with the wrong name, or a .pptx that is really a legacy
        # .ppt. python-pptx says "File is not a zip file", which tells a
        # lecturer nothing. This function's contract is a message the UI can
        # show, so translate rather than propagate.
        raise ValueError(
            f"{filename} could not be read as a {suffix.upper()} file. It may be "
            "corrupt, incomplete, or saved in a different format than its name "
            "suggests."
        ) from exc

    logger.info(
        "Parsed %s: %d blocks, %d containing maths",
        filename, len(blocks), sum(1 for b in blocks if b.get("has_math")),
    )
    return blocks


def blocks_needing_math_ocr(blocks: List[Block]) -> List[Block]:
    """PDF blocks that look mathematical but whose maths did not survive extraction."""
    return [b for b in blocks if b.get("needs_math_ocr")]
